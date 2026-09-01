"""Local PowerPoint authoring and editing for Echo."""

from __future__ import annotations

import contextlib
from pathlib import Path
from typing import Any

from runtime.execution.suckers.registry import Skill
from runtime.platform.plugins.bundled._office_io import (
    atomic_package_save,
    create_versioned_backup,
    replace_text_preserving_runs,
    scoped_path_denial,
)
from runtime.platform.plugins.plugin_base import ModulePlugin

PLUGIN_NAME = "presentations"
_TRUSTED_SOURCE = "plugin://presentations"

try:  # pragma: no cover - dependency probe
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    _PPTX_OK = True
except Exception:  # pragma: no cover
    _PPTX_OK = False


def _require_pptx() -> dict[str, Any] | None:
    if _PPTX_OK:
        return None
    return {
        "ok": False,
        "error": "python-pptx 未安装，无法处理 pptx（安装: pip install python-pptx）",
    }


def _resolve_pptx(path: Any, *, write: bool = False) -> tuple[Path | None, dict[str, Any] | None]:
    if not isinstance(path, (str, Path)) or not str(path).strip():
        return None, {"ok": False, "error": "path 不能为空"}
    try:
        resolved = Path(str(path)).expanduser().resolve()
    except Exception as exc:
        return None, {"ok": False, "error": f"无效路径: {exc}"}
    if resolved.suffix.lower() != ".pptx":
        return None, {"ok": False, "error": "仅支持 .pptx 文件"}
    denial = scoped_path_denial(resolved, write=write)
    if denial:
        return None, {"ok": False, "error": denial}
    return resolved, None


def _iter_text_frames(slide: Any):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape.text_frame
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


class PresentationsPlugin(ModulePlugin):
    name = PLUGIN_NAME
    display_name = "Presentations"
    version = "0.1.0"
    description = "本地创建、提取和原位修改 PowerPoint 演示文稿。"
    author = "Echo"

    def register_skills(self) -> None:
        if self.ctx is None:
            return
        skills = [
            Skill(
                name="presentations.create_pptx",
                description=(
                    "从结构化 slides 创建 16:9 pptx。参数:path，slides:[{title,subtitle,"
                    "bullets,kind}]，overwrite 默认 false。"
                ),
                summary="创建 pptx(path+slides)",
                affinity=[
                    "presentations",
                    "pptx",
                    "powerpoint",
                    "slides",
                    "file",
                    "write",
                    "create",
                ],
                cost_profile="low",
                trusted_source=_TRUSTED_SOURCE,
                handler=self._create_pptx,
            ),
            Skill(
                name="presentations.extract_text",
                description="按页提取 pptx 的形状、表格文本及备注。参数:path。",
                summary="提取 pptx 文本",
                affinity=["presentations", "pptx", "powerpoint", "file", "read"],
                cost_profile="low",
                trusted_source=_TRUSTED_SOURCE,
                handler=self._extract_text,
            ),
            Skill(
                name="presentations.replace_text",
                description=(
                    "在现有 pptx 中精确替换文本，可用 slide 限定页码，并保留其他版式。"
                    "参数:path，replacements:[{old,new}]，slide 可选(1-based)，backup 默认 true。"
                ),
                summary="原位修改 pptx 文本",
                affinity=[
                    "presentations",
                    "pptx",
                    "powerpoint",
                    "file",
                    "write",
                    "edit",
                    "replace",
                ],
                cost_profile="low",
                trusted_source=_TRUSTED_SOURCE,
                handler=self._replace_text,
            ),
            Skill(
                name="presentations.presentation_info",
                description="返回 pptx 页数、尺寸、形状、表格和图片统计。参数:path。",
                summary="pptx 结构统计",
                affinity=["presentations", "pptx", "powerpoint", "file", "read", "info"],
                cost_profile="low",
                trusted_source=_TRUSTED_SOURCE,
                handler=self._presentation_info,
            ),
        ]
        for skill in skills:
            with contextlib.suppress(Exception):
                self.ctx.register_skill(skill)

    def _create_pptx(self, **kwargs: Any) -> dict[str, Any]:
        path, err = _resolve_pptx(kwargs.get("path"), write=True)
        if err:
            return err
        assert path is not None
        if path.exists() and not kwargs.get("overwrite"):
            return {"ok": False, "error": f"目标已存在: {path},需 overwrite=true"}
        dependency_error = _require_pptx()
        if dependency_error:
            return dependency_error
        slides = kwargs.get("slides")
        if not isinstance(slides, list) or not slides:
            return {"ok": False, "error": "slides 必须是非空数组"}

        deck = Presentation()
        deck.slide_width = Inches(13.333333)
        deck.slide_height = Inches(7.5)
        for index, spec in enumerate(slides):
            if not isinstance(spec, dict):
                return {"ok": False, "error": f"slides[{index}] 必须是对象"}
            kind = str(spec.get("kind") or ("title" if index == 0 else "content"))
            layout_index = 0 if kind == "title" else 1
            layout_index = min(layout_index, len(deck.slide_layouts) - 1)
            slide = deck.slides.add_slide(deck.slide_layouts[layout_index])
            if slide.shapes.title is not None:
                slide.shapes.title.text = str(spec.get("title") or "")
            subtitle = str(spec.get("subtitle") or "")
            bullets = spec.get("bullets") or []
            if not isinstance(bullets, list):
                return {"ok": False, "error": f"slides[{index}].bullets 必须是数组"}
            body = next(
                (
                    shape
                    for shape in slide.placeholders
                    if shape != slide.shapes.title and getattr(shape, "has_text_frame", False)
                ),
                None,
            )
            if body is not None:
                frame = body.text_frame
                frame.clear()
                content = [subtitle] if subtitle else []
                content.extend(str(item) for item in bullets)
                for item_index, text in enumerate(content):
                    paragraph = frame.paragraphs[0] if item_index == 0 else frame.add_paragraph()
                    paragraph.text = text
                    paragraph.level = 0
                    paragraph.font.size = Pt(24 if kind == "title" else 20)
                    paragraph.font.color.rgb = RGBColor(45, 52, 65)
            if slide.shapes.title is not None:
                for paragraph in slide.shapes.title.text_frame.paragraphs:
                    paragraph.font.size = Pt(34 if kind == "title" else 30)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(22, 28, 45)
        try:
            atomic_package_save(path, deck.save)
        except Exception as exc:  # noqa: BLE001
            return {"ok": False, "error": f"保存 pptx 失败: {exc}"}
        return {"ok": True, "path": str(path), "slides": len(deck.slides)}

    def _open(self, path_value: Any):
        path, err = _resolve_pptx(path_value)
        if err:
            return None, None, err
        assert path is not None
        if not path.is_file():
            return None, None, {"ok": False, "error": f"文件不存在: {path}"}
        dependency_error = _require_pptx()
        if dependency_error:
            return None, None, dependency_error
        try:
            return path, Presentation(path), None
        except Exception as exc:  # noqa: BLE001
            return None, None, {"ok": False, "error": f"无法打开 pptx: {exc}"}

    def _extract_text(self, **kwargs: Any) -> dict[str, Any]:
        path, deck, err = self._open(kwargs.get("path"))
        if err:
            return err
        slides: list[dict[str, Any]] = []
        for slide_index, slide in enumerate(deck.slides, start=1):
            blocks: list[dict[str, Any]] = []
            for shape_index, shape in enumerate(slide.shapes, start=1):
                if getattr(shape, "has_table", False):
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    blocks.append(
                        {
                            "shape": shape_index,
                            "name": str(shape.name),
                            "type": "table",
                            "rows": rows,
                            "text": "\n".join(" | ".join(row) for row in rows),
                        }
                    )
                    continue
                text = str(getattr(shape, "text", "") or "").strip()
                if text:
                    blocks.append(
                        {
                            "shape": shape_index,
                            "name": str(shape.name),
                            "type": "text",
                            "text": text,
                        }
                    )
            notes = ""
            with contextlib.suppress(Exception):
                notes = str(slide.notes_slide.notes_text_frame.text or "").strip()
            slides.append({"slide": slide_index, "blocks": blocks, "notes": notes})
        return {"ok": True, "path": str(path), "slides": slides}

    def _replace_text(self, **kwargs: Any) -> dict[str, Any]:
        path, deck, err = self._open(kwargs.get("path"))
        if err:
            return err
        assert path is not None
        denial = scoped_path_denial(path, write=True)
        if denial:
            return {"ok": False, "error": denial}
        raw = kwargs.get("replacements")
        if not isinstance(raw, list) or not raw:
            return {"ok": False, "error": "replacements 必须是非空 [{old,new}] 数组"}
        replacements: list[tuple[str, str]] = []
        for index, item in enumerate(raw):
            if not isinstance(item, dict) or not str(item.get("old") or ""):
                return {"ok": False, "error": f"replacements[{index}].old 无效"}
            replacements.append((str(item["old"]), str(item.get("new") or "")))
        try:
            requested_slide = int(kwargs.get("slide") or 0)
        except (TypeError, ValueError):
            return {"ok": False, "error": "slide 必须是整数页码"}
        if requested_slide < 0 or requested_slide > len(deck.slides):
            return {"ok": False, "error": f"页码超出范围: {requested_slide}"}
        counts = [0] * len(replacements)
        target_slides = [deck.slides[requested_slide - 1]] if requested_slide else list(deck.slides)
        for slide in target_slides:
            for frame in _iter_text_frames(slide):
                for paragraph in frame.paragraphs:
                    replace_text_preserving_runs(paragraph, replacements, counts)
        total = sum(counts)
        if total == 0:
            return {"ok": False, "error": "未找到待替换文本，文件未改动"}
        backup_path = create_versioned_backup(path) if kwargs.get("backup", True) else None
        try:
            atomic_package_save(path, deck.save)
        except Exception as exc:  # noqa: BLE001
            return {"ok": False, "error": f"保存 pptx 失败: {exc}"}
        return {
            "ok": True,
            "path": str(path),
            "backup_path": str(backup_path) if backup_path else None,
            "replacement_counts": counts,
            "total_replacements": total,
        }

    def _presentation_info(self, **kwargs: Any) -> dict[str, Any]:
        path, deck, err = self._open(kwargs.get("path"))
        if err:
            return err
        slides = []
        for index, slide in enumerate(deck.slides, start=1):
            shapes = list(slide.shapes)
            slides.append(
                {
                    "slide": index,
                    "shapes": len(shapes),
                    "tables": sum(1 for shape in shapes if getattr(shape, "has_table", False)),
                    "pictures": sum(
                        1 for shape in shapes if int(getattr(shape, "shape_type", 0)) == 13
                    ),
                }
            )
        return {
            "ok": True,
            "path": str(path),
            "slide_count": len(deck.slides),
            "width": int(deck.slide_width),
            "height": int(deck.slide_height),
            "slides": slides,
            "size_bytes": path.stat().st_size,
        }


__all__ = ["PresentationsPlugin"]
